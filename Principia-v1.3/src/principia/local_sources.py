from __future__ import annotations

import csv
import hashlib
import io
import json
import mimetypes
import os
import re
import unicodedata
import xml.etree.ElementTree as ET
from collections.abc import Callable, Iterator
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from urllib.parse import quote

from pypdf import PdfReader

from .ids import short_hash
from .models import (
    LocalCorpusConfig,
    LocalCorpusDiagnostics,
    LocalSourceReport,
    WorkItem,
    WorkList,
)
from .storage import WorkspaceStorage

ARCHIVE_EXTENSIONS = {
    ".7z",
    ".bz2",
    ".gz",
    ".rar",
    ".tar",
    ".tgz",
    ".xz",
    ".zip",
}
LEGACY_OFFICE_EXTENSIONS = {".doc", ".ppt", ".xls"}
OPAQUE_MEDIA_EXTENSIONS = {
    ".aac",
    ".avi",
    ".flac",
    ".gif",
    ".jpeg",
    ".jpg",
    ".m4a",
    ".mkv",
    ".mov",
    ".mp3",
    ".mp4",
    ".mpeg",
    ".png",
    ".tiff",
    ".wav",
    ".webm",
    ".webp",
}
TEXT_EXTENSIONS = {
    ".bib",
    ".c",
    ".cc",
    ".cfg",
    ".conf",
    ".cpp",
    ".css",
    ".f",
    ".f90",
    ".go",
    ".h",
    ".hpp",
    ".ini",
    ".java",
    ".jl",
    ".js",
    ".kt",
    ".log",
    ".m",
    ".md",
    ".py",
    ".r",
    ".rb",
    ".rst",
    ".rs",
    ".sh",
    ".sql",
    ".tex",
    ".toml",
    ".ts",
    ".txt",
    ".yaml",
    ".yml",
}


class LocalParserUnavailable(RuntimeError):
    """Raised when an optional parser dependency is not installed."""


@dataclass(frozen=True)
class ParsedLocalContent:
    text: str
    content_type: str = "local_text"
    warnings: tuple[str, ...] = ()


LocalParser = Callable[[Path, bytes], str | ParsedLocalContent]


@dataclass(frozen=True)
class RegisteredLocalParser:
    name: str
    parser: LocalParser
    extensions: tuple[str, ...]
    mime_types: tuple[str, ...]
    version: str
    fingerprint: str


class LocalParserRegistry:
    """Extension/MIME parser registry with stable cache fingerprints."""

    def __init__(self) -> None:
        self._parsers: list[RegisteredLocalParser] = []

    def register(
        self,
        name: str,
        parser: LocalParser,
        *,
        extensions: tuple[str, ...] | list[str] = (),
        mime_types: tuple[str, ...] | list[str] = (),
        version: str = "1",
    ) -> RegisteredLocalParser:
        normalized_extensions = tuple(
            sorted(
                {
                    item.lower() if item.startswith(".") else f".{item.lower()}"
                    for item in extensions
                }
            )
        )
        normalized_mimes = tuple(sorted({item.lower() for item in mime_types}))
        identity = {
            "name": name,
            "version": str(version),
            "extensions": normalized_extensions,
            "mime_types": normalized_mimes,
            "callable": f"{getattr(parser, '__module__', '')}.{getattr(parser, '__qualname__', repr(parser))}",
        }
        registered = RegisteredLocalParser(
            name=str(name),
            parser=parser,
            extensions=normalized_extensions,
            mime_types=normalized_mimes,
            version=str(version),
            fingerprint=hashlib.sha256(
                json.dumps(identity, sort_keys=True).encode("utf-8")
            ).hexdigest(),
        )
        self._parsers = [item for item in self._parsers if item.name != registered.name]
        self._parsers.append(registered)
        return registered

    def resolve(self, path: Path, mime_type: str) -> RegisteredLocalParser | None:
        suffix = path.suffix.lower()
        for parser in reversed(self._parsers):
            if suffix in parser.extensions:
                return parser
        normalized_mime = mime_type.lower()
        for parser in reversed(self._parsers):
            if normalized_mime and normalized_mime in parser.mime_types:
                return parser
        return None


DEFAULT_LOCAL_PARSERS = LocalParserRegistry()


def register_local_parser(
    name: str,
    parser: LocalParser,
    *,
    extensions: tuple[str, ...] | list[str] = (),
    mime_types: tuple[str, ...] | list[str] = (),
    version: str = "1",
) -> RegisteredLocalParser:
    """Register an OCR, EPUB, transcription, or organization-specific parser."""

    return DEFAULT_LOCAL_PARSERS.register(
        name,
        parser,
        extensions=extensions,
        mime_types=mime_types,
        version=version,
    )


def chunk_local_text(text: str, *, chunk_chars: int = 24_000, overlap: int = 2_000) -> list[str]:
    """Split normalized text deterministically with bounded overlap."""

    if chunk_chars < 1 or overlap < 0 or overlap >= chunk_chars:
        raise ValueError("Require chunk_chars > overlap >= 0")
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        hard_end = min(len(text), start + chunk_chars)
        end = hard_end
        if hard_end < len(text):
            boundary = max(
                text.rfind("\n", start + chunk_chars // 2, hard_end),
                text.rfind(" ", start + chunk_chars // 2, hard_end),
            )
            if boundary > start:
                end = boundary
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        next_start = max(start + 1, end - overlap)
        start = next_start
    return chunks


class LocalCorpusIngestor:
    def __init__(
        self,
        storage: WorkspaceStorage,
        *,
        registry: LocalParserRegistry | None = None,
    ) -> None:
        self.storage = storage
        self.registry = registry or DEFAULT_LOCAL_PARSERS

    def ingest(
        self,
        folder: str | Path,
        *,
        config: LocalCorpusConfig | None = None,
    ) -> WorkList:
        config = config or LocalCorpusConfig()
        root = Path(folder).expanduser()
        if not root.exists():
            raise FileNotFoundError(f"Local corpus folder does not exist: {root}")
        if not root.is_dir():
            raise NotADirectoryError(f"Local corpus path is not a directory: {root}")
        root = root.resolve()
        corpus_name = _corpus_slug(config.corpus_name or root.name)
        diagnostics = LocalCorpusDiagnostics(corpus_name=corpus_name)
        works: list[WorkItem] = []
        seen_hashes: dict[str, str] = {}
        candidates = list(
            _iter_corpus_entries(
                root,
                recursive=config.recursive,
                include_hidden=config.include_hidden,
            )
        )
        diagnostics.discovered_count = len(candidates)
        if len(candidates) > config.max_files:
            diagnostics.warnings.append(
                f"Corpus contained {len(candidates)} entries; only the first {config.max_files} were inspected."
            )
            candidates = candidates[: config.max_files]

        for path in candidates:
            report, work = self._ingest_path(root, path, corpus_name, config, seen_hashes)
            diagnostics.reports.append(report)
            if report.status == "accepted":
                diagnostics.accepted_count += 1
            elif report.status == "cached":
                diagnostics.cached_count += 1
            elif report.status == "duplicate":
                diagnostics.duplicate_count += 1
            elif report.status == "error":
                diagnostics.failed_count += 1
            else:
                diagnostics.skipped_count += 1
            diagnostics.total_bytes += report.byte_size
            diagnostics.total_characters += report.character_count
            if work is not None:
                works.append(work)

        return WorkList(
            query=f"local corpus:{corpus_name}",
            items=works,
            target_count=0,
            mode="local",
            sources=["local"],
            local_diagnostics=diagnostics,
        )

    def _ingest_path(
        self,
        root: Path,
        path: Path,
        corpus_name: str,
        config: LocalCorpusConfig,
        seen_hashes: dict[str, str],
    ) -> tuple[LocalSourceReport, WorkItem | None]:
        relative = path.relative_to(root).as_posix()
        uri = _portable_local_uri(corpus_name, relative)
        base: dict[str, Any] = {"uri": uri, "relative_path": relative}
        if path.is_symlink():
            return (
                LocalSourceReport(
                    **base,
                    status="skipped",
                    warnings=["Symbolic links are not ingested."],
                ),
                None,
            )
        try:
            stat = path.stat()
        except OSError as exc:
            return LocalSourceReport(**base, status="error", error=_safe_error(exc)), None
        if not path.is_file():
            return (
                LocalSourceReport(
                    **base,
                    status="skipped",
                    warnings=["Non-regular filesystem entries are not ingested."],
                ),
                None,
            )
        suffix = path.suffix.lower()
        if suffix in ARCHIVE_EXTENSIONS:
            return _skipped_report(
                base, stat.st_size, "Archives are not expanded or ingested."
            ), None
        if suffix in LEGACY_OFFICE_EXTENSIONS:
            return _skipped_report(
                base,
                stat.st_size,
                "Legacy Office files are unsupported; save as DOCX, PPTX, or XLSX.",
            ), None
        if suffix in OPAQUE_MEDIA_EXTENSIONS:
            return _skipped_report(
                base,
                stat.st_size,
                "Opaque media requires a registered OCR or transcription parser.",
            ), None
        if stat.st_size > config.max_file_bytes:
            return (
                _skipped_report(
                    base,
                    stat.st_size,
                    f"File exceeds the configured {config.max_file_bytes}-byte limit.",
                ),
                None,
            )
        try:
            data = path.read_bytes()
        except OSError as exc:
            return LocalSourceReport(
                **base, status="error", byte_size=stat.st_size, error=_safe_error(exc)
            ), None
        byte_hash = hashlib.sha256(data).hexdigest()
        mime_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
        duplicate_uri = seen_hashes.get(byte_hash)
        if duplicate_uri and duplicate_uri != uri:
            return (
                LocalSourceReport(
                    **base,
                    status="duplicate",
                    mime_type=mime_type,
                    byte_sha256=byte_hash,
                    byte_size=len(data),
                    duplicate_of=duplicate_uri,
                    warnings=["Exact-byte duplicate skipped."],
                ),
                None,
            )
        stored_duplicate = self.storage.get_source_asset_by_byte_hash(byte_hash)
        if stored_duplicate and stored_duplicate.get("portable_uri") != uri:
            duplicate_of = str(stored_duplicate.get("portable_uri") or "existing local source")
            seen_hashes[byte_hash] = duplicate_of
            return (
                LocalSourceReport(
                    **base,
                    status="duplicate",
                    mime_type=mime_type,
                    byte_sha256=byte_hash,
                    byte_size=len(data),
                    duplicate_of=duplicate_of,
                    warnings=["Exact-byte duplicate already exists in this workspace."],
                ),
                None,
            )
        parser = self.registry.resolve(path, mime_type)
        if parser is None:
            parser = _unknown_text_parser_registration(path, mime_type)
        existing = self.storage.get_source_asset_by_uri(uri)
        if (
            existing
            and existing.get("byte_sha256") == byte_hash
            and existing.get("parser_fingerprint") == parser.fingerprint
            and existing.get("normalized_text")
        ):
            work = self.storage.get_work(str(existing["work_id"]))
            if work is not None:
                seen_hashes[byte_hash] = uri
                return _report_from_asset(existing, status="cached"), work
        try:
            parsed_value = parser.parser(path, data)
            parsed = (
                parsed_value
                if isinstance(parsed_value, ParsedLocalContent)
                else ParsedLocalContent(str(parsed_value))
            )
            text = normalize_local_text(parsed.text)
            if not text:
                raise ValueError("Parser returned no usable text")
        except LocalParserUnavailable as exc:
            return (
                LocalSourceReport(
                    **base,
                    status="skipped",
                    mime_type=mime_type,
                    parser=parser.name,
                    parser_fingerprint=parser.fingerprint,
                    byte_sha256=byte_hash,
                    byte_size=len(data),
                    warnings=[str(exc)],
                ),
                None,
            )
        except Exception as exc:  # noqa: BLE001
            return (
                LocalSourceReport(
                    **base,
                    status="error",
                    mime_type=mime_type,
                    parser=parser.name,
                    parser_fingerprint=parser.fingerprint,
                    byte_sha256=byte_hash,
                    byte_size=len(data),
                    error=_safe_error(exc),
                ),
                None,
            )

        text_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()
        chunks = chunk_local_text(
            text,
            chunk_chars=config.chunk_chars,
            overlap=config.chunk_overlap,
        )
        chunk_hashes = [hashlib.sha256(item.encode("utf-8")).hexdigest() for item in chunks]
        work = WorkItem(
            id=f"L-{short_hash(uri, length=20).upper()}",
            title=_document_title(path),
            source="local",
            source_type="local_document",
            url=uri,
            source_urls=[uri],
            content_sha256=byte_hash,
            metadata={
                "private": True,
                "local_uri": uri,
                "corpus_name": corpus_name,
                "mime_type": mime_type,
                "parser": parser.name,
                "parser_fingerprint": parser.fingerprint,
                "text_sha256": text_hash,
                "character_count": len(text),
                "chunk_count": len(chunks),
                "chunk_chars": config.chunk_chars,
                "chunk_overlap": config.chunk_overlap,
            },
        )
        work = self.storage.save_work(work)
        asset_id = f"A-{short_hash(uri, length=20).upper()}"
        asset = self.storage.save_source_asset(
            {
                "id": asset_id,
                "work_id": work.id,
                "corpus_name": corpus_name,
                "portable_uri": uri,
                "relative_path": relative,
                "absolute_path": str(path.resolve()),
                "mime_type": mime_type,
                "parser_name": parser.name,
                "parser_fingerprint": parser.fingerprint,
                "byte_sha256": byte_hash,
                "text_sha256": text_hash,
                "byte_size": len(data),
                "character_count": len(text),
                "chunk_count": len(chunks),
                "chunk_hashes": chunk_hashes,
                "chunk_chars": config.chunk_chars,
                "chunk_overlap": config.chunk_overlap,
                "content_type": parsed.content_type,
                "normalized_text": text,
                "status": "accepted",
                "warnings": list(parsed.warnings),
            }
        )
        seen_hashes[byte_hash] = uri
        return _report_from_asset(asset, status="accepted"), work


def normalize_local_text(text: str) -> str:
    value = unicodedata.normalize("NFKC", str(text or ""))
    value = value.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[\t ]+", " ", line).strip() for line in value.split("\n")]
    output: list[str] = []
    blank = False
    for line in lines:
        if not line:
            if output and not blank:
                output.append("")
            blank = True
            continue
        output.append(line)
        blank = False
    return "\n".join(output).strip()


def _iter_corpus_entries(root: Path, *, recursive: bool, include_hidden: bool) -> Iterator[Path]:
    def visit(directory: Path) -> Iterator[Path]:
        try:
            entries = sorted(os.scandir(directory), key=lambda item: item.name.casefold())
        except OSError:
            return
        for entry in entries:
            if not include_hidden and entry.name.startswith("."):
                continue
            path = Path(entry.path)
            if entry.is_symlink():
                yield path
            elif entry.is_dir(follow_symlinks=False):
                if recursive:
                    yield from visit(path)
            else:
                yield path

    yield from visit(root)


class _VisibleHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._suppressed = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() in {"script", "style", "noscript"}:
            self._suppressed += 1
        elif tag.lower() in {"br", "p", "div", "h1", "h2", "h3", "li", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {"script", "style", "noscript"} and self._suppressed:
            self._suppressed -= 1
        elif tag.lower() in {"p", "div", "h1", "h2", "h3", "li", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._suppressed:
            self.parts.append(data)


def _parse_pdf(path: Path, data: bytes) -> ParsedLocalContent:
    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted and reader.decrypt("") == 0:
        raise ValueError("Encrypted PDF requires a password and was not ingested")
    pages = [str(page.extract_text() or "") for page in reader.pages]
    return ParsedLocalContent("\n\n".join(pages), content_type="pdf_text")


def _parse_html(path: Path, data: bytes) -> ParsedLocalContent:
    parser = _VisibleHTMLParser()
    parser.feed(_decode_text(data))
    return ParsedLocalContent("".join(parser.parts), content_type="html")


def _parse_xml(path: Path, data: bytes) -> ParsedLocalContent:
    root = ET.fromstring(data)
    return ParsedLocalContent("\n".join(item.strip() for item in root.itertext() if item.strip()))


def _parse_json(path: Path, data: bytes) -> ParsedLocalContent:
    value = json.loads(_decode_text(data))
    return ParsedLocalContent(json.dumps(value, ensure_ascii=False, indent=2, sort_keys=True))


def _parse_jsonl(path: Path, data: bytes) -> ParsedLocalContent:
    values = [json.loads(line) for line in _decode_text(data).splitlines() if line.strip()]
    return ParsedLocalContent(
        "\n".join(json.dumps(item, ensure_ascii=False, sort_keys=True) for item in values)
    )


def _parse_csv(path: Path, data: bytes) -> ParsedLocalContent:
    text = _decode_text(data)
    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
    except csv.Error:
        dialect = csv.excel_tab if path.suffix.lower() == ".tsv" else csv.excel
    rows = list(csv.reader(io.StringIO(text), dialect=dialect))
    return ParsedLocalContent("\n".join(" | ".join(cell.strip() for cell in row) for row in rows))


def _parse_text(path: Path, data: bytes) -> ParsedLocalContent:
    return ParsedLocalContent(_decode_text(data))


def _parse_docx(path: Path, data: bytes) -> ParsedLocalContent:
    try:
        from docx import Document  # type: ignore[import-not-found, import-untyped]
    except ImportError as exc:  # pragma: no cover - depends on optional extra
        raise LocalParserUnavailable(
            "DOCX parsing requires `pip install principia-ai[local]`."
        ) from exc
    document = Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        parts.extend(" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows)
    return ParsedLocalContent("\n".join(parts))


def _parse_pptx(path: Path, data: bytes) -> ParsedLocalContent:
    try:
        from pptx import Presentation  # type: ignore[import-not-found, import-untyped]
    except ImportError as exc:  # pragma: no cover - depends on optional extra
        raise LocalParserUnavailable(
            "PPTX parsing requires `pip install principia-ai[local]`."
        ) from exc
    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = [
            str(shape.text)
            for shape in slide.shapes
            if hasattr(shape, "text") and str(shape.text).strip()
        ]
        if slide_text:
            parts.append(f"Slide {index}\n" + "\n".join(slide_text))
    return ParsedLocalContent("\n\n".join(parts))


def _parse_xlsx(path: Path, data: bytes) -> ParsedLocalContent:
    try:
        from openpyxl import load_workbook  # type: ignore[import-untyped]
    except ImportError as exc:  # pragma: no cover - depends on optional extra
        raise LocalParserUnavailable(
            "XLSX parsing requires `pip install principia-ai[local]`."
        ) from exc
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [
                    str(value).strip() for value in row if value is not None and str(value).strip()
                ]
                if values:
                    parts.append(" | ".join(values))
    finally:
        workbook.close()
    return ParsedLocalContent("\n".join(parts))


def _decode_text(data: bytes) -> str:
    if not data:
        return ""
    encodings = ["utf-8-sig"]
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        encodings.insert(0, "utf-16")
    for encoding in encodings:
        try:
            text = data.decode(encoding)
        except UnicodeDecodeError:
            continue
        if _printable_ratio(text) >= 0.85:
            return text
    text = data.decode("latin-1")
    if _printable_ratio(text) < 0.85:
        raise ValueError("File is not safely decodable text; register a format-specific parser")
    return text


def _printable_ratio(text: str) -> float:
    if not text:
        return 1.0
    printable = sum(character.isprintable() or character in "\n\r\t" for character in text)
    return printable / len(text)


def _unknown_text_parser_registration(path: Path, mime_type: str) -> RegisteredLocalParser:
    identity = {
        "name": "safe_text",
        "version": "1",
        "mime_type": mime_type,
        "suffix": path.suffix.lower(),
    }
    return RegisteredLocalParser(
        name="safe_text",
        parser=_parse_text,
        extensions=(path.suffix.lower(),) if path.suffix else (),
        mime_types=(mime_type,),
        version="1",
        fingerprint=hashlib.sha256(
            json.dumps(identity, sort_keys=True).encode("utf-8")
        ).hexdigest(),
    )


def _document_title(path: Path) -> str:
    title = re.sub(r"[_-]+", " ", path.stem).strip()
    return title or path.name or "Local document"


def _corpus_slug(value: str) -> str:
    normalized = unicodedata.normalize("NFKC", str(value or "corpus")).strip()
    normalized = re.sub(r"[^\w.-]+", "-", normalized, flags=re.UNICODE).strip("-.")
    return normalized or "corpus"


def _portable_local_uri(corpus_name: str, relative_path: str) -> str:
    return f"local://{quote(corpus_name, safe='')}/{quote(relative_path, safe='/')}"


def _skipped_report(base: dict[str, Any], byte_size: int, warning: str) -> LocalSourceReport:
    return LocalSourceReport.model_validate(
        {**base, "status": "skipped", "byte_size": byte_size, "warnings": [warning]}
    )


def _report_from_asset(asset: dict[str, Any], *, status: str) -> LocalSourceReport:
    return LocalSourceReport(
        uri=str(asset.get("portable_uri") or ""),
        relative_path=str(asset.get("relative_path") or ""),
        status=status,  # type: ignore[arg-type]
        work_id=str(asset.get("work_id") or ""),
        mime_type=str(asset.get("mime_type") or ""),
        parser=str(asset.get("parser_name") or ""),
        parser_fingerprint=str(asset.get("parser_fingerprint") or ""),
        byte_sha256=str(asset.get("byte_sha256") or ""),
        text_sha256=str(asset.get("text_sha256") or ""),
        byte_size=int(asset.get("byte_size") or 0),
        character_count=int(asset.get("character_count") or 0),
        chunk_count=int(asset.get("chunk_count") or 0),
        warnings=[str(item) for item in asset.get("warnings") or []],
    )


def _safe_error(exc: Exception) -> str:
    detail = str(exc)[:300]
    detail = re.sub(r"(?<!\w)(?:[A-Za-z]:[\\/]|/)[^\s,;]+", "[local path]", detail)
    return f"{type(exc).__name__}: {detail}"


DEFAULT_LOCAL_PARSERS.register(
    "pdf",
    _parse_pdf,
    extensions=(".pdf",),
    mime_types=("application/pdf",),
    version="pypdf-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "html",
    _parse_html,
    extensions=(".html", ".htm"),
    mime_types=("text/html",),
    version="htmlparser-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "xml",
    _parse_xml,
    extensions=(".xml",),
    mime_types=("application/xml", "text/xml"),
    version="elementtree-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "json",
    _parse_json,
    extensions=(".json",),
    mime_types=("application/json",),
    version="json-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "jsonl", _parse_jsonl, extensions=(".jsonl", ".ndjson"), version="jsonl-v1"
)
DEFAULT_LOCAL_PARSERS.register(
    "tabular_text",
    _parse_csv,
    extensions=(".csv", ".tsv"),
    mime_types=("text/csv", "text/tab-separated-values"),
    version="csv-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "text",
    _parse_text,
    extensions=tuple(TEXT_EXTENSIONS),
    mime_types=("text/plain", "text/markdown", "text/x-python", "application/x-yaml"),
    version="text-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "docx",
    _parse_docx,
    extensions=(".docx",),
    mime_types=("application/vnd.openxmlformats-officedocument.wordprocessingml.document",),
    version="python-docx-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "pptx",
    _parse_pptx,
    extensions=(".pptx",),
    mime_types=("application/vnd.openxmlformats-officedocument.presentationml.presentation",),
    version="python-pptx-v1",
)
DEFAULT_LOCAL_PARSERS.register(
    "xlsx",
    _parse_xlsx,
    extensions=(".xlsx",),
    mime_types=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",),
    version="openpyxl-v1",
)
